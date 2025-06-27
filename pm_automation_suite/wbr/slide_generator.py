"""
Slide Generator Implementation

Automated slide generation for WBR presentations with template engine,
dynamic chart creation, brand compliance, and multi-format support.
"""

import logging
import json
import os
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass
from enum import Enum
import base64
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # Create dummy classes for type hints when pptx is not available
    class Presentation:
        pass

try:
    import matplotlib.pyplot as plt
    import seaborn as sns
    import pandas as pd
    PLOTTING_AVAILABLE = True
except ImportError:
    PLOTTING_AVAILABLE = False

try:
    from connectors.google_connector import GoogleConnector
    GOOGLE_AVAILABLE = True
except ImportError:
    GOOGLE_AVAILABLE = False

# Handle missing modules for type hints
if not PPTX_AVAILABLE:
    class Inches:
        @staticmethod
        def __call__(val):
            return val
    
    class Pt:
        @staticmethod
        def __call__(val):
            return val

from .wbr_data_extractor import WBRDataPackage, WBRMetric
from .insight_generator import Insight


logger = logging.getLogger(__name__)


class SlideType(Enum):
    """Types of slides."""
    TITLE = "title"
    EXECUTIVE_SUMMARY = "executive_summary"
    METRICS_OVERVIEW = "metrics_overview"
    TREND_ANALYSIS = "trend_analysis"
    DEEP_DIVE = "deep_dive"
    INSIGHTS = "insights"
    ACTION_ITEMS = "action_items"
    APPENDIX = "appendix"


class ChartType(Enum):
    """Types of charts."""
    LINE = "line"
    BAR = "bar"
    PIE = "pie"
    SCATTER = "scatter"
    HEATMAP = "heatmap"
    GAUGE = "gauge"


@dataclass
class SlideTemplate:
    """Template configuration for slides."""
    slide_type: SlideType
    title: str
    layout_name: str
    content_placeholders: List[str]
    chart_config: Optional[Dict[str, Any]] = None
    style_overrides: Optional[Dict[str, Any]] = None


@dataclass
class ChartConfig:
    """Configuration for chart generation."""
    chart_type: ChartType
    title: str
    data_source: str
    x_axis: str
    y_axis: str
    series: List[str]
    colors: Optional[List[str]] = None
    style: Optional[Dict[str, Any]] = None


class SlideGenerator:
    """
    Automated slide generator for WBR presentations.
    
    Creates PowerPoint presentations with dynamic charts, insights,
    and branded templates based on WBR data packages.
    """
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize Slide Generator.
        
        Args:
            config: Configuration dictionary with template and brand settings
        """
        self.config = config
        self.brand_config = config.get('brand', {})
        self.template_path = config.get('template_path')
        self.output_path = config.get('output_path', './wbr_presentations')
        
        # Brand colors (default corporate palette)
        self.brand_colors = {
            'primary': self.brand_config.get('primary_color', '#1f4e79'),
            'secondary': self.brand_config.get('secondary_color', '#70ad47'),
            'accent': self.brand_config.get('accent_color', '#ffc000'),
            'warning': self.brand_config.get('warning_color', '#ff0000'),
            'text': self.brand_config.get('text_color', '#333333'),
            'background': self.brand_config.get('background_color', '#ffffff')
        }
        
        # Chart styling
        self.chart_style = {
            'font_family': self.brand_config.get('font_family', 'Calibri'),
            'font_size': self.brand_config.get('font_size', 12),
            'grid_alpha': 0.3,
            'figure_size': (10, 6),
            'dpi': 300
        }
        
        # Load slide templates
        self.slide_templates = self._load_slide_templates()
        
        # Initialize Google Connector if available
        self.google_connector = None
        if GOOGLE_AVAILABLE and 'google' in config:
            self.google_connector = GoogleConnector(config['google'])
        
        # Ensure output directory exists
        os.makedirs(self.output_path, exist_ok=True)
        
    def _load_slide_templates(self) -> Dict[SlideType, SlideTemplate]:
        """Load slide templates configuration."""
        templates = {
            SlideType.TITLE: SlideTemplate(
                slide_type=SlideType.TITLE,
                title="Weekly Business Review",
                layout_name="Title Slide",
                content_placeholders=["title", "subtitle", "date"],
                style_overrides={'title_size': 44, 'subtitle_size': 24}
            ),
            
            SlideType.EXECUTIVE_SUMMARY: SlideTemplate(
                slide_type=SlideType.EXECUTIVE_SUMMARY,
                title="Executive Summary",
                layout_name="Title and Content",
                content_placeholders=["summary_text", "key_metrics"],
                style_overrides={'bullet_level': 1}
            ),
            
            SlideType.METRICS_OVERVIEW: SlideTemplate(
                slide_type=SlideType.METRICS_OVERVIEW,
                title="Key Metrics Overview",
                layout_name="Title and Content",
                content_placeholders=["metrics_table", "trend_chart"],
                chart_config={
                    'type': ChartType.BAR,
                    'show_targets': True,
                    'show_previous': True
                }
            ),
            
            SlideType.TREND_ANALYSIS: SlideTemplate(
                slide_type=SlideType.TREND_ANALYSIS,
                title="Trend Analysis",
                layout_name="Title and Content",
                content_placeholders=["trend_chart", "insights"],
                chart_config={
                    'type': ChartType.LINE,
                    'show_trend_lines': True,
                    'highlight_changes': True
                }
            ),
            
            SlideType.INSIGHTS: SlideTemplate(
                slide_type=SlideType.INSIGHTS,
                title="Key Insights & Recommendations",
                layout_name="Title and Content",
                content_placeholders=["insights_list", "priority_actions"],
                style_overrides={'highlight_critical': True}
            ),
            
            SlideType.ACTION_ITEMS: SlideTemplate(
                slide_type=SlideType.ACTION_ITEMS,
                title="Action Items",
                layout_name="Title and Content",
                content_placeholders=["action_table", "owners", "timeline"],
                style_overrides={'table_style': 'branded'}
            )
        }
        
        return templates
    
    async def generate_wbr_presentation(
        self, 
        wbr_data: WBRDataPackage,
        insights: List[Insight],
        executive_summary: str,
        custom_slides: Optional[List[Dict[str, Any]]] = None
    ) -> str:
        """
        Generate complete WBR presentation.
        
        Args:
            wbr_data: Complete WBR data package
            insights: Generated insights
            executive_summary: Executive summary text
            custom_slides: Optional custom slide configurations
            
        Returns:
            Path to generated presentation file
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available. Cannot generate presentations.")
            raise ImportError("python-pptx package required for slide generation")
        
        logger.info("Starting WBR presentation generation")
        
        try:
            # Create presentation
            prs = self._create_presentation()
            
            # Generate slides in order
            self._add_title_slide(prs, wbr_data.generation_time)
            self._add_executive_summary_slide(prs, executive_summary, wbr_data)
            self._add_metrics_overview_slide(prs, wbr_data.metrics)
            self._add_trend_analysis_slide(prs, wbr_data.metrics, wbr_data.raw_data)
            self._add_insights_slide(prs, insights)
            self._add_action_items_slide(prs, insights)
            
            # Add custom slides if provided
            if custom_slides:
                for slide_config in custom_slides:
                    self._add_custom_slide(prs, slide_config)
            
            # Add appendix with raw data
            self._add_appendix_slides(prs, wbr_data)
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"WBR_{timestamp}.pptx"
            filepath = os.path.join(self.output_path, filename)
            
            prs.save(filepath)
            
            logger.info(f"WBR presentation saved to: {filepath}")
            
            # Upload to Google Drive if configured
            if self.google_connector:
                try:
                    drive_url = await self._upload_to_google_drive(filepath, filename)
                    logger.info(f"Presentation uploaded to Google Drive: {drive_url}")
                except Exception as e:
                    logger.warning(f"Failed to upload to Google Drive: {e}")
            
            return filepath
            
        except Exception as e:
            logger.error(f"Failed to generate WBR presentation: {e}")
            raise
    
    def _create_presentation(self) -> Presentation:
        """Create base presentation with template."""
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
            logger.info(f"Using template: {self.template_path}")
        else:
            prs = Presentation()
            logger.info("Using default presentation template")
        
        return prs
    
    def _add_title_slide(self, prs: Presentation, generation_time: datetime):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Weekly Business Review"
        
        # Set subtitle with date
        subtitle = slide.placeholders[1]
        week_start = generation_time.strftime("%B %d")
        week_end = (generation_time).strftime("%B %d, %Y")
        subtitle.text = f"Week of {week_start} - {week_end}"
        
        # Apply brand styling
        self._apply_brand_styling(title.text_frame, 'title')
        self._apply_brand_styling(subtitle.text_frame, 'subtitle')
    
    def _add_executive_summary_slide(
        self, 
        prs: Presentation, 
        executive_summary: str,
        wbr_data: WBRDataPackage
    ):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Executive Summary"
        self._apply_brand_styling(title.text_frame, 'title')
        
        # Add executive summary text
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = executive_summary
        
        # Add data quality indicator
        p = text_frame.add_paragraph()
        p.text = f"\\nData Quality Score: {wbr_data.quality_score:.1f}/1.0"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['text']))
        
        self._apply_brand_styling(text_frame, 'body')
    
    def _add_metrics_overview_slide(self, prs: Presentation, metrics: List[WBRMetric]):
        """Add metrics overview slide with table and chart."""
        slide_layout = prs.slide_layouts[5]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Key Metrics Overview"
        self._apply_brand_styling(title.text_frame, 'title')
        
        # Create metrics table
        self._add_metrics_table(slide, metrics)
        
        # Add trend chart
        if PLOTTING_AVAILABLE:
            chart_image = self._create_metrics_chart(metrics)
            if chart_image:
                self._add_image_to_slide(slide, chart_image, left=Inches(5), top=Inches(2), width=Inches(5))
    
    def _add_trend_analysis_slide(
        self, 
        prs: Presentation, 
        metrics: List[WBRMetric],
        raw_data: Dict[str, Any]
    ):
        """Add trend analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Trend Analysis"
        self._apply_brand_styling(title.text_frame, 'title')
        
        # Create trend chart
        if PLOTTING_AVAILABLE and raw_data:
            chart_image = self._create_trend_chart(raw_data)
            if chart_image:
                content = slide.placeholders[1]
                # Add chart
                self._add_image_to_slide(slide, chart_image, left=Inches(1), top=Inches(1.5), width=Inches(8))
                
                # Add trend insights
                insights_text = self._generate_trend_insights(metrics)
                if insights_text:
                    text_frame = content.text_frame
                    text_frame.text = insights_text
                    self._apply_brand_styling(text_frame, 'body')
    
    def _add_insights_slide(self, prs: Presentation, insights: List[Insight]):
        """Add insights slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Key Insights & Recommendations"
        self._apply_brand_styling(title.text_frame, 'title')
        
        # Add insights content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Group insights by priority
        critical_insights = [i for i in insights if i.priority.value == 'critical']
        high_insights = [i for i in insights if i.priority.value == 'high']
        
        if critical_insights:
            p = text_frame.paragraphs[0]
            p.text = "🚨 Critical Issues:"
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['warning']))
            
            for insight in critical_insights[:2]:  # Limit to top 2
                p = text_frame.add_paragraph()
                p.text = f"• {insight.title}: {insight.description}"
                p.level = 1
        
        if high_insights:
            p = text_frame.add_paragraph()
            p.text = "\\n⚡ High Priority:"
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['secondary']))
            
            for insight in high_insights[:3]:  # Limit to top 3
                p = text_frame.add_paragraph()
                p.text = f"• {insight.title}: {insight.description}"
                p.level = 1
        
        self._apply_brand_styling(text_frame, 'body')
    
    def _add_action_items_slide(self, prs: Presentation, insights: List[Insight]):
        """Add action items slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Action Items"
        self._apply_brand_styling(title.text_frame, 'title')
        
        # Create action items table
        self._add_action_items_table(slide, insights)
    
    def _add_custom_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """Add custom slide based on configuration."""
        slide_layout = prs.slide_layouts[slide_config.get('layout', 1)]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if 'title' in slide_config:
            title = slide.shapes.title
            title.text = slide_config['title']
            self._apply_brand_styling(title.text_frame, 'title')
        
        # Add content based on type
        if slide_config.get('type') == 'chart':
            # Add chart content
            pass  # Implementation would depend on chart configuration
        elif slide_config.get('type') == 'text':
            # Add text content
            if 'content' in slide_config:
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.text = slide_config['content']
                self._apply_brand_styling(text_frame, 'body')
    
    def _add_appendix_slides(self, prs: Presentation, wbr_data: WBRDataPackage):
        """Add appendix slides with detailed data."""
        # Data sources slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Appendix: Data Sources"
        self._apply_brand_styling(title.text_frame, 'title')
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        # Add data freshness information
        freshness_info = []
        for source, timestamp in wbr_data.data_freshness.items():
            freshness_info.append(f"• {source}: {timestamp.strftime('%Y-%m-%d %H:%M')}")
        
        text_frame.text = "Data Freshness:\\n" + "\\n".join(freshness_info)
        text_frame.add_paragraph().text = f"\\nGenerated: {wbr_data.generation_time.strftime('%Y-%m-%d %H:%M:%S')}"
        
        self._apply_brand_styling(text_frame, 'body')
    
    def _add_metrics_table(self, slide, metrics: List[WBRMetric]):
        """Add metrics table to slide."""
        # Table dimensions
        rows = len(metrics) + 1  # +1 for header
        cols = 5  # Metric, Current, Previous, Change, Target
        
        # Add table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set header
        headers = ["Metric", "Current", "Previous", "Change", "Target"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['primary']))
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Fill data rows
        for i, metric in enumerate(metrics, 1):
            table.cell(i, 0).text = metric.name
            table.cell(i, 1).text = f"{metric.value:.2f}"
            table.cell(i, 2).text = f"{metric.previous_value:.2f}" if metric.previous_value else "N/A"
            
            # Format change cell with color
            change_cell = table.cell(i, 3)
            if metric.change_percent is not None:
                change_text = f"{metric.change_percent:+.1f}%"
                change_cell.text = change_text
                
                # Color based on trend
                if metric.trend == "up":
                    change_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['secondary']))
                elif metric.trend == "down":
                    change_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['warning']))
            else:
                change_cell.text = "N/A"
            
            table.cell(i, 4).text = f"{metric.target:.2f}" if metric.target else "N/A"
    
    def _add_action_items_table(self, slide, insights: List[Insight]):
        """Add action items table to slide."""
        # Collect all recommendations
        action_items = []
        for insight in insights:
            for rec in insight.recommendations[:2]:  # Limit recommendations per insight
                action_items.append({
                    'action': rec,
                    'priority': insight.priority.value,
                    'source': insight.title
                })
        
        if not action_items:
            return
        
        # Limit to top 8 action items
        action_items = action_items[:8]
        
        # Table dimensions
        rows = len(action_items) + 1
        cols = 3  # Action, Priority, Owner
        
        # Add table
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.4 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set header
        headers = ["Action Item", "Priority", "Owner"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['primary']))
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Fill data rows
        for i, item in enumerate(action_items, 1):
            table.cell(i, 0).text = item['action']
            
            # Priority cell with color
            priority_cell = table.cell(i, 1)
            priority_cell.text = item['priority'].title()
            if item['priority'] == 'critical':
                priority_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['warning']))
            elif item['priority'] == 'high':
                priority_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['secondary']))
            
            table.cell(i, 2).text = "TBD"  # Owner to be determined
    
    def _create_metrics_chart(self, metrics: List[WBRMetric]) -> Optional[BytesIO]:
        """Create metrics chart."""
        if not PLOTTING_AVAILABLE:
            return None
        
        try:
            # Prepare data
            metric_names = [m.name[:20] + "..." if len(m.name) > 20 else m.name for m in metrics]
            current_values = [m.value for m in metrics]
            targets = [m.target if m.target else 0 for m in metrics]
            
            # Create chart
            fig, ax = plt.subplots(figsize=self.chart_style['figure_size'], 
                                 dpi=self.chart_style['dpi'])
            
            x_pos = range(len(metric_names))
            
            # Current values bars
            bars1 = ax.bar([x - 0.2 for x in x_pos], current_values, 
                          width=0.4, label='Current', 
                          color=self.brand_colors['primary'], alpha=0.8)
            
            # Target values bars
            bars2 = ax.bar([x + 0.2 for x in x_pos], targets, 
                          width=0.4, label='Target', 
                          color=self.brand_colors['secondary'], alpha=0.8)
            
            # Styling
            ax.set_xlabel('Metrics')
            ax.set_ylabel('Values')
            ax.set_title('Current vs Target Performance')
            ax.set_xticks(x_pos)
            ax.set_xticklabels(metric_names, rotation=45, ha='right')
            ax.legend()
            ax.grid(True, alpha=self.chart_style['grid_alpha'])
            
            # Apply brand styling
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            plt.tight_layout()
            
            # Save to BytesIO
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            return img_buffer
            
        except Exception as e:
            logger.error(f"Failed to create metrics chart: {e}")
            return None
    
    def _create_trend_chart(self, raw_data: Dict[str, Any]) -> Optional[BytesIO]:
        """Create trend chart from raw data."""
        if not PLOTTING_AVAILABLE or not raw_data:
            return None
        
        try:
            # Use user activity data if available
            if 'user_activity' in raw_data:
                df = raw_data['user_activity']
                
                if not df.empty and 'date' in df.columns:
                    fig, ax = plt.subplots(figsize=self.chart_style['figure_size'],
                                         dpi=self.chart_style['dpi'])
                    
                    # Plot trends
                    if 'unique_users' in df.columns:
                        ax.plot(df['date'], df['unique_users'], 
                               marker='o', color=self.brand_colors['primary'],
                               linewidth=2, label='Daily Active Users')
                    
                    if 'total_events' in df.columns:
                        ax2 = ax.twinx()
                        ax2.plot(df['date'], df['total_events'],
                                marker='s', color=self.brand_colors['secondary'],
                                linewidth=2, label='Total Events')
                        ax2.set_ylabel('Total Events')
                        ax2.legend(loc='upper right')
                    
                    # Styling
                    ax.set_xlabel('Date')
                    ax.set_ylabel('Daily Active Users')
                    ax.set_title('User Activity Trends')
                    ax.legend(loc='upper left')
                    ax.grid(True, alpha=self.chart_style['grid_alpha'])
                    
                    # Rotate x-axis labels
                    plt.xticks(rotation=45)
                    plt.tight_layout()
                    
                    # Save to BytesIO
                    img_buffer = BytesIO()
                    plt.savefig(img_buffer, format='png', bbox_inches='tight')
                    img_buffer.seek(0)
                    plt.close()
                    
                    return img_buffer
            
        except Exception as e:
            logger.error(f"Failed to create trend chart: {e}")
            
        return None
    
    def _add_image_to_slide(self, slide, image_buffer: BytesIO, left, top, width):
        """Add image to slide from BytesIO buffer."""
        try:
            slide.shapes.add_picture(image_buffer, left, top, width=width)
        except Exception as e:
            logger.error(f"Failed to add image to slide: {e}")
    
    def _generate_trend_insights(self, metrics: List[WBRMetric]) -> str:
        """Generate text insights about trends."""
        insights = []
        
        improving_metrics = [m for m in metrics if m.trend == "up"]
        declining_metrics = [m for m in metrics if m.trend == "down"]
        
        if improving_metrics:
            insights.append(f"📈 Improving: {', '.join([m.name for m in improving_metrics[:3]])}")
        
        if declining_metrics:
            insights.append(f"📉 Declining: {', '.join([m.name for m in declining_metrics[:3]])}")
        
        return "\\n".join(insights)
    
    def _apply_brand_styling(self, text_frame, style_type: str):
        """Apply brand styling to text frame."""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.chart_style['font_family']
                    
                    if style_type == 'title':
                        run.font.size = Pt(28)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['primary']))
                    elif style_type == 'subtitle':
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['text']))
                    elif style_type == 'body':
                        run.font.size = Pt(self.chart_style['font_size'])
                        run.font.color.rgb = RGBColor(*self._hex_to_rgb(self.brand_colors['text']))
                        
        except Exception as e:
            logger.error(f"Failed to apply brand styling: {e}")
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    async def _upload_to_google_drive(self, filepath: str, filename: str) -> Optional[str]:
        """Upload presentation to Google Drive."""
        if not self.google_connector:
            return None
        
        try:
            # Upload file to Google Drive
            file_metadata = {
                'name': filename,
                'parents': [self.config.get('google_drive_folder_id')]
            }
            
            # This would be implemented based on Google Connector capabilities
            # drive_url = await self.google_connector.upload_file(filepath, file_metadata)
            # return drive_url
            
            logger.info("Google Drive upload would happen here")
            return None
            
        except Exception as e:
            logger.error(f"Google Drive upload failed: {e}")
            return None
    
    def get_slide_preview(self, slide_type: SlideType) -> str:
        """Get preview/description of slide type."""
        previews = {
            SlideType.TITLE: "Title slide with WBR branding and date range",
            SlideType.EXECUTIVE_SUMMARY: "High-level summary with key metrics and data quality",
            SlideType.METRICS_OVERVIEW: "Table and chart showing current vs target performance",
            SlideType.TREND_ANALYSIS: "Line charts showing metric trends over time",
            SlideType.INSIGHTS: "AI-generated insights organized by priority",
            SlideType.ACTION_ITEMS: "Table of recommended actions with priorities",
            SlideType.APPENDIX: "Data sources, freshness, and technical details"
        }
        
        return previews.get(slide_type, "Custom slide content")
    
    async def generate_slide_preview(self, slide_type: SlideType, data: Dict[str, Any]) -> str:
        """Generate HTML preview of slide content."""
        # This would generate an HTML preview of the slide
        # Implementation would depend on requirements
        return f"<div>Preview of {slide_type.value} slide</div>"